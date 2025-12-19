# app.py – Klaidų analizė su automatine PPT, HH:MM:SS laiku, finansine rizika ir grafika
import streamlit as st
import pandas as pd
import plotly.express as px
from pptx import Presentation
from pptx.util import Inches
import io
import os

st.set_page_config(page_title="Klaidų analizė", layout="wide")
st.title("📊 Klaidų analizė su PowerPoint generavimu")

# ----------------------------
# 1. DUOMENŲ ĮKĖLIMAS
# ----------------------------
uploaded_file = st.file_uploader("Įkelkite Excel klaidų registrą", type=["xlsx"])

if uploaded_file:
    df = pd.read_excel(uploaded_file)

    # ----------------------------
    # 2. KLAIDOS IDENTIFIKACIJA
    # ----------------------------
    df["Yra klaida"] = df["Klaidos tipas"].notna() & (df["Klaidos tipas"].astype(str).str.strip() != "")
    klaidos_df = df[df["Yra klaida"]].copy()

    # ----------------------------
    # 3. FINANSINĖ RIZIKA
    # ----------------------------
    def nustatyti_finansine_rizika(row):
        try:
            suma = float(row.get("Suma EUR, be PVM", 0))
        except:
            suma = 0
        klaidos_tipas = str(row.get("Klaidos tipas", "")).lower()
        if "terminas" in klaidos_tipas:
            return suma
        else:
            return 0

    klaidos_df["Finansinė rizika (€)"] = klaidos_df.apply(nustatyti_finansine_rizika, axis=1)

    # ----------------------------
    # 4. TAISYMO LAIKAS
    # ----------------------------
    klaidos_df["Pradžia"] = pd.to_datetime(
        klaidos_df["Klaidos ištaisymo laiko pradžia"], format="%H:%M:%S", errors="coerce"
    )
    klaidos_df["Pabaiga"] = pd.to_datetime(
        klaidos_df["Klaidos ištaisymo laiko pabaiga"], format="%H:%M:%S", errors="coerce"
    )

    klaidos_df["Taisymo laikas (min)"] = (
        (klaidos_df["Pabaiga"] - klaidos_df["Pradžia"]).dt.total_seconds() / 60
    )
    klaidos_df["Taisymo laikas (min)"] = klaidos_df["Taisymo laikas (min)"].apply(
        lambda x: x + 24*60 if x < 0 else x
    )
    klaidos_df["Taisymo laikas (val)"] = klaidos_df["Taisymo laikas (min)"] / 60

    # ----------------------------
    # 5. KLAIDOS SUNKUMAS
    # ----------------------------
    def nustatyti_sunkuma(row):
        rizika = row.get("Finansinė rizika (€)", 0)
        laikas = row.get("Taisymo laikas (min)", 0)
        if rizika >= 1000 or laikas >= 240:
            return "Kritinė"
        elif rizika >= 100 or laikas >= 60:
            return "Vidutinė"
        elif rizika > 0:
            return "Maža"
        else:
            return "Administracinė"

    klaidos_df["Klaidos sunkumas"] = klaidos_df.apply(nustatyti_sunkuma, axis=1)

    # ----------------------------
    # 6. KPI
    # ----------------------------
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("📌 Tikrų klaidų skaičius", len(klaidos_df))
    col2.metric("⏱️ Prarastas laikas (val)", round(klaidos_df["Taisymo laikas (val)"].sum(), 2))
    col3.metric("💰 Bendra finansinė rizika (€)", round(klaidos_df["Finansinė rizika (€)"].sum(), 2))
    col4.metric("🔥 Kritinių klaidų", (klaidos_df["Klaidos sunkumas"] == "Kritinė").sum())

    # ----------------------------
    # 7. GRAFIKAI
    # ----------------------------
    st.subheader("📈 Klaidų pasiskirstymas pagal sunkumą")
    fig1 = px.bar(
        klaidos_df.groupby("Klaidos sunkumas").size().reset_index(name="Kiekis"),
        x="Klaidos sunkumas",
        y="Kiekis",
        color="Klaidos sunkumas"
    )
    st.plotly_chart(fig1, use_container_width=True)

    st.subheader("🏭 Klaidos pagal proceso etapą")
    fig2 = px.bar(
        klaidos_df.groupby("Proceso etapas").size().reset_index(name="Kiekis"),
        x="Proceso etapas",
        y="Kiekis"
    )
    st.plotly_chart(fig2, use_container_width=True)

    st.subheader("⏱️ Klaidų taisymo laikas (val)")
    fig3 = px.bar(
        klaidos_df,
        x="Klaidos tipas",
        y="Taisymo laikas (val)",
        color="Klaidos sunkumas",
        hover_data=["Finansinė rizika (€)", "Proceso etapas", "Atsakinga puse"]
    )
    st.plotly_chart(fig3, use_container_width=True)

    # ----------------------------
    # 8. AUTOMATINĖ POWERPOINT GENERACIJA SU GRAFIKA
    # ----------------------------
    if st.button("📤 Generuoti PowerPoint su grafika"):
        prs = Presentation()
        layout_blank = prs.slide_layouts[5]

        # KPI slide
        slide = prs.slides.add_slide(layout_blank)
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(2))
        tf = txBox.text_frame
        tf.text = f"Tikrų klaidų skaičius: {len(klaidos_df)}\n"
        tf.add_paragraph().text = f"Prarastas laikas (val): {round(klaidos_df['Taisymo laikas (val)'].sum(), 2)}"
        tf.add_paragraph().text = f"Bendra finansinė rizika (€): {round(klaidos_df['Finansinė rizika (€)'].sum(),2)}"
        tf.add_paragraph().text = f"Kritinių klaidų skaičius: {(klaidos_df['Klaidos sunkumas'] == 'Kritinė').sum()}"

        # Grafikai – eksportuojame į PNG
        os.makedirs("temp_figs", exist_ok=True)
        fig1.write_image("temp_figs/fig1.png")
        fig2.write_image("temp_figs/fig2.png")
        fig3.write_image("temp_figs/fig3.png")

        # Įdedame grafikus į slide
        for fig_file, title in zip(["temp_figs/fig1.png","temp_figs/fig2.png","temp_figs/fig3.png"],
                                   ["Klaidų pasiskirstymas pagal sunkumą","Klaidos pagal proceso etapą","Klaidų taisymo laikas"]):
            slide = prs.slides.add_slide(layout_blank)
            slide.shapes.add_picture(fig_file, Inches(0.5), Inches(1), width=Inches(9))
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.5))
            txBox.text = title

        # Visos klaidos – batch po 15 eilučių
        batch_size = 15
        for i in range(0, len(klaidos_df), batch_size):
            slide = prs.slides.add_slide(layout_blank)
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(6))
            tf = txBox.text_frame
            for j in range(i, min(i+batch_size, len(klaidos_df))):
                row = klaidos_df.iloc[j]
                p = tf.add_paragraph()
                p.text = f"{row['Klaidos tipas']} | {row['Finansinė rizika (€)']} € | {row['Taisymo laikas (val)']:.2f} val | {row['Proceso etapas']} | {row['Atsakinga puse']}"

        # Save į in-memory
        pptx_io = io.BytesIO()
        prs.save(pptx_io)
        pptx_io.seek(0)

        st.success("PowerPoint su grafika paruošta!")
        st.download_button("📥 Atsisiųsti PowerPoint", pptx_io, file_name="Klaidu_ataskaita.pptx")

else:
    st.info("Įkelkite Excel failą, kad pradėtume analizę")
